from pathlib import Path
from fastapi import UploadFile
from langchain.schema import SystemMessage, HumanMessage
from PIL import Image
import re
from io import BytesIO
from typing import List, Dict, Tuple
from sklearn.metrics.pairwise import cosine_similarity
from google.cloud import texttospeech_v1 as tts
from pptx import Presentation
from pptx.util import Inches
from typing import Optional
import tempfile, zipfile
import base64
import fitz
import io
import os
import uuid

PDF_DIR = Path(r"..\data\save_pdf")
txt_DIR = Path(r"..\data\save_txt")
IMAGE_DIR = Path(r"..\data\temp_images")
AUDIO_DIR = Path("../data/audio")

def preprocess_text(text: str, max_length: int = 2000) -> str:
    """텍스트 전처리 및 길이 제한"""
    # 불필요한 공백 제거
    text = re.sub(r'\s+', ' ', text)
    # 특수문자 정리
    text = re.sub(r'[^\w\s.,!?-]', '', text)
    # 길이 제한
    return text[:max_length]

def convert_image_to_base64(image_bytes: bytes) -> str:
    """
    이미지 바이트를 base64 문자열로 변환
    """
    return base64.b64encode(image_bytes).decode('utf-8')


def optimize_image(image_bytes: bytes, max_size: int = 800, quality: int = 50) -> bytes:
    """
    이미지 크기 리사이즈 + JPEG 압축
    """
    img = Image.open(BytesIO(image_bytes)).convert("RGB")
    if img.width > max_size or img.height > max_size:
        ratio = max_size / max(img.width, img.height)
        new_size = (int(img.width * ratio), int(img.height * ratio))
        img = img.resize(new_size, Image.Resampling.LANCZOS)
    
    output = BytesIO()
    img.save(output, format="JPEG", quality=quality, optimize=True)
    return output.getvalue(), img.size

def extract_image_bytes(doc, xref):
    """
    PDF에서 이미지 바이트 데이터와 사이즈를 추출
    """
    pix = fitz.Pixmap(doc, xref)
    if pix.n > 4:  # CMYK → RGB 변환
        pix = fitz.Pixmap(fitz.csRGB, pix)
    img_bytes = pix.tobytes()
    
    return optimize_image(img_bytes)

def extract_page_bytes(page):
    """
    PDF에서 페이지 바이트 데이터와 사이즈를 추출
    """
    pix = page.get_pixmap(matrix=fitz.Matrix(150/72, 150/72))  # 150 DPI
    page_bytes = pix.tobytes("png")

    return optimize_image(page_bytes)

def calculate_image_ratio(image_size, page_size):
    """
    이미지가 페이지에서 차지하는 비율 계산
    """
    image_area = image_size[0] * image_size[1]
    page_area = page_size[0] * page_size[1]

    return image_area / page_area

def clear_audio_dir(audio_dir: Path):
    """기존 음성 파일 모두 삭제"""
    for file in audio_dir.glob("*.wav"):
        file.unlink()

def export_pdf_with_audio_to_pptx(pdf_bytes: bytes, wav_dir: str) -> bytes:
    """
    PDF 파일을 PPTX로 변환하고, 각 페이지에 대응되는 오디오(WAV)를 삽입하여 반환.

    Parameters:
    - pdf_bytes: Streamlit에서 업로드한 PDF의 byte stream
    - wav_dir: 페이지별 WAV 파일이 저장된 디렉터리 (파일명: page_0.wav, page_1.wav, ...)

    Returns:
    - PPTX byte stream (다운로드용)
    """
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    ppt = Presentation()

    for page_index in range(len(doc)):
        page = doc.load_page(page_index)

        # 📷 PDF 페이지 → 이미지로 렌더링
        image_path = os.path.join(tempfile.gettempdir(), f"slide_{page_index}.png")
        page.get_pixmap(matrix=fitz.Matrix(2, 2)).save(image_path)

        # 🎞️ 새 슬라이드 추가 + 이미지 삽입
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

        # 🔈 오디오 삽입
        audio_path = os.path.join(wav_dir, f"page_{page_index}.wav")
        if os.path.exists(audio_path):
            try:
                slide.shapes.add_movie(
                    audio_path,
                    left=Inches(0.5), top=Inches(0.5),
                    width=Inches(1), height=Inches(1),
                    mime_type='audio/wav'
                )
                print(f"✅ Slide {page_index+1}: 오디오 삽입 완료")
            except Exception as e:
                print(f"⚠️ Slide {page_index+1}: 오디오 삽입 실패 → {e}")
        else:
            print(f"⚠️ Slide {page_index+1}: WAV 파일 없음 → {audio_path}")

    # 💾 결과 임시 파일로 저장 후 byte 반환
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        ppt.save(tmp.name)
        tmp.seek(0)
        return tmp.read()
    
def export_pptx_with_wavs_as_zip(pptx_bytes: bytes, wav_dir: str) -> bytes:
    """PPTX 파일과 WAV 파일들을 하나의 ZIP으로 묶어 반환"""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".zip") as tmp_zip:
        with zipfile.ZipFile(tmp_zip.name, 'w') as zipf:
            # PPTX 파일 추가
            zipf.writestr("presentation.pptx", pptx_bytes)
            
            # WAV 파일들 추가
            for filename in os.listdir(wav_dir):
                if filename.endswith(".wav"):
                    filepath = os.path.join(wav_dir, filename)
                    zipf.write(filepath, arcname=os.path.join("audio", filename))

        tmp_zip.seek(0)
        return tmp_zip.read()